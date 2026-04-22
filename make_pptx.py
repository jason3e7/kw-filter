#!/usr/bin/env python3
"""Generate kw-filter.pptx — improved visual layout."""
from __future__ import annotations
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

# ── Palette ───────────────────────────────────────────────────────────────────
BG       = RGBColor(0x0f, 0x17, 0x2a)
BG2      = RGBColor(0x1e, 0x29, 0x3b)
BG3      = RGBColor(0x09, 0x11, 0x1e)   # code bg
ACCENT   = RGBColor(0x4a, 0xde, 0x80)
DANGER   = RGBColor(0xf8, 0x71, 0x71)
WARN     = RGBColor(0xfb, 0xbf, 0x24)
BLUE     = RGBColor(0x60, 0xa5, 0xfa)
PURPLE   = RGBColor(0xc0, 0x84, 0xfc)
MUTED    = RGBColor(0x94, 0xa3, 0xb8)
WHITE    = RGBColor(0xf1, 0xf5, 0xf9)
BORDER   = RGBColor(0x33, 0x41, 0x55)
CODE_FG  = RGBColor(0xa5, 0xf3, 0xfc)

# ── Layout constants ──────────────────────────────────────────────────────────
W  = Inches(13.33)
H  = Inches(7.5)
ML = Inches(0.65)           # margin left
MT = Inches(0.45)           # margin top (heading)
CW = Inches(12.03)          # content width
HW = Inches(5.86)           # half-width column
GAP= Inches(0.31)           # column gap
R  = ML + HW + GAP          # right column x
BODY_TOP = Inches(1.18)     # content top (below heading)
BODY_H   = Inches(5.9)      # available content height
CODE_PT  = Pt(10.5)         # code font size — fits ~65 chars in HW col


# ── Helpers ───────────────────────────────────────────────────────────────────

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return slide


def _set_txbody_margins(txBody, l=Pt(8), r=Pt(8), t=Pt(6), b=Pt(6), anchor='t'):
    bp = txBody.find(qn('a:bodyPr'))
    if bp is not None:
        bp.set('lIns', str(int(l)))
        bp.set('rIns', str(int(r)))
        bp.set('tIns', str(int(t)))
        bp.set('bIns', str(int(b)))
        bp.set('anchor', anchor)


def _set_anchor(txBody, anchor='t'):
    bp = txBody.find(qn('a:bodyPr'))
    if bp is not None:
        bp.set('anchor', anchor)


def txtbox(slide, x, y, w, h, text, size=Pt(14), color=WHITE,
           bold=False, italic=False, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = "Noto Sans TC"
    return tb


def heading(slide, text, y=MT):
    """Slide heading + green underline rule."""
    txtbox(slide, ML, y, CW, Inches(0.62), text,
           size=Pt(30), color=WHITE, bold=True)
    rule = slide.shapes.add_shape(1, ML, y + Inches(0.63), CW, Pt(2.5))
    rule.fill.solid(); rule.fill.fore_color.rgb = ACCENT
    rule.line.fill.background()


def rect(slide, x, y, w, h, fill_color=BG2, border_color=BORDER, border_pt=1.2):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    return shape


def card(slide, x, y, w, h, title="", body="",
         border=BORDER, title_color=ACCENT, body_color=WHITE,
         body_size=Pt(13.5), fill=BG2):
    shape = rect(slide, x, y, w, h, fill_color=fill, border_color=border)
    tf = shape.text_frame
    tf.word_wrap = True
    _set_txbody_margins(tf._txBody, l=Pt(11), r=Pt(11), t=Pt(9), b=Pt(9))
    if title:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(15)
        run.font.bold = True
        run.font.color.rgb = title_color
        run.font.name = "Noto Sans TC"
    if body:
        p2 = tf.paragraphs[0] if not title else tf.add_paragraph()
        if title:
            p2.space_before = Pt(4)
        p2.alignment = PP_ALIGN.LEFT
        run2 = p2.add_run()
        run2.text = body
        run2.font.size = body_size
        run2.font.color.rgb = body_color
        run2.font.name = "Noto Sans TC"
    return shape


def code_block(slide, x, y, w, h, code: str, label="", label_color=MUTED):
    if label:
        txtbox(slide, x, y, w, Inches(0.3), label,
               size=Pt(12), color=label_color)
        y += Inches(0.3); h -= Inches(0.3)
    shape = rect(slide, x, y, w, h, fill_color=BG3, border_color=BORDER, border_pt=1)
    tf = shape.text_frame
    tf.word_wrap = False
    _set_txbody_margins(tf._txBody, l=Pt(11), r=Pt(11), t=Pt(9), b=Pt(9))
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = Pt(17)
        run = p.add_run()
        run.text = line
        run.font.size = CODE_PT
        run.font.color.rgb = CODE_FG
        run.font.name = "Courier New"
    return shape


def table(slide, x, y, w, col_fracs: list,
          headers: list, rows: list,
          row_h=Inches(0.44), cell_size=Pt(13)):
    n_cols = len(headers)
    n_rows = len(rows) + 1
    tbl = slide.shapes.add_table(n_rows, n_cols, x, y, w, row_h * n_rows)

    total_fracs = sum(col_fracs)
    for ci, frac in enumerate(col_fracs):
        tbl.table.columns[ci].width = int(w * frac / total_fracs)

    def fmt(cell, text, bold=False, fg=WHITE, bg=BG2, size=cell_size, align=PP_ALIGN.LEFT):
        cell.fill.solid(); cell.fill.fore_color.rgb = bg
        tf = cell.text_frame; tf.word_wrap = True
        _set_txbody_margins(tf._txBody, l=Pt(7), r=Pt(5), t=Pt(5), b=Pt(5))
        p = tf.paragraphs[0]; p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = size; run.font.bold = bold
        run.font.color.rgb = fg; run.font.name = "Noto Sans TC"

    HDR_BG = RGBColor(0x0d, 0x24, 0x18)
    for ci, h in enumerate(headers):
        fmt(tbl.table.cell(0, ci), h, bold=True, fg=ACCENT, bg=HDR_BG)

    for ri, row in enumerate(rows):
        bg = RGBColor(0x14, 0x1e, 0x32) if ri % 2 == 0 else BG
        for ci, val in enumerate(row):
            fg = RGBColor(0x6e, 0xe7, 0xb7) if ci == 0 else WHITE
            fmt(tbl.table.cell(ri + 1, ci), str(val), fg=fg, bg=bg)

    return tbl


def stat_box(slide, x, y, w, h, number, label):
    shape = rect(slide, x, y, w, h, fill_color=BG2, border_color=BORDER)
    tf = shape.text_frame; tf.word_wrap = True
    _set_anchor(tf._txBody, 'ctr')
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = number; r.font.size = Pt(40); r.font.bold = True
    r.font.color.rgb = ACCENT; r.font.name = "Noto Sans TC"
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label; r2.font.size = Pt(13)
    r2.font.color.rgb = MUTED; r2.font.name = "Noto Sans TC"


def flow_row(slide, items: list[tuple], x, y, h=Inches(0.5)):
    """items = [(label, is_cmd, color)]  arrows auto-inserted between"""
    cx = x
    for i, (label, is_cmd, color) in enumerate(items):
        is_arr = label == "→"
        fw = Inches(0.3) if is_arr else Inches(1.15)
        if is_arr:
            shape = rect(slide, cx, y, fw, h, fill_color=BG, border_color=None)
        else:
            bg = RGBColor(0x0e, 0x2a, 0x1a) if is_cmd else BG2
            bc = ACCENT if is_cmd else BORDER
            shape = rect(slide, cx, y, fw, h, fill_color=bg, border_color=bc)
        tf = shape.text_frame
        _set_anchor(tf._txBody, 'ctr')
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.size = Pt(12.5)
        run.font.bold = is_cmd
        run.font.color.rgb = color if not is_arr else MUTED
        run.font.name = "Noto Sans TC"
        cx += fw + Inches(0.055)


# ══════════════════════════════════════════════════════════════════════════════
# Slides
# ══════════════════════════════════════════════════════════════════════════════

def s01_cover(prs):
    slide = blank(prs)

    # ── Vertically center the whole block ──
    # Block: badge(0.3) + gap(0.25) + title(1.05) + sub(0.5) + desc(0.45) + gap(0.5) + stats(1.55) = 4.6
    START = (7.5 - 4.6) / 2  # ≈ 1.45
    START = Inches(START)

    txtbox(slide, ML, START, CW, Inches(0.32),
           "Open Source  ·  Python  ·  Zero Dependencies",
           size=Pt(13.5), color=MUTED, align=PP_ALIGN.CENTER)

    txtbox(slide, ML, START + Inches(0.38), CW, Inches(1.05),
           "kw-filter",
           size=Pt(62), color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

    txtbox(slide, ML, START + Inches(1.5), CW, Inches(0.52),
           "保護隱私的 AI 資料過濾工具",
           size=Pt(24), color=WHITE, align=PP_ALIGN.CENTER)

    txtbox(slide, ML, START + Inches(2.08), CW, Inches(0.45),
           "在資料送給 AI 之前，先把敏感關鍵字換掉；拿回答案之後，再還原回來。",
           size=Pt(14.5), color=MUTED, align=PP_ALIGN.CENTER)

    # Stats — 4 equal boxes
    sw = Inches(2.75); sh = Inches(1.6)
    sx = ML + Inches(0.26)
    sy = START + Inches(2.75)
    for num, lbl in [("6","指令"), ("160+","自動化測試"), ("0","第三方依賴"), ("29萬+","行/秒處理")]:
        stat_box(slide, sx, sy, sw, sh, num, lbl)
        sx += sw + Inches(0.29)


def s02_problem(prs):
    slide = blank(prs)
    heading(slide, "你每天都在洩漏什麼？")

    txtbox(slide, ML, BODY_TOP, CW, Inches(0.32),
           "把文件直接貼給 AI 前，想一下裡面有什麼……",
           size=Pt(14), color=MUTED)

    cw = Inches(3.9); cy = BODY_TOP + Inches(0.4); ch = Inches(2.1)
    for title, bc, body in [
        ("👤 個人資料", DANGER, "姓名、Email\n電話、身分證\n客戶資料、合約對象"),
        ("🔑 機密憑證", DANGER, "API Key / Token\n資料庫密碼\nPrivate Key 路徑"),
        ("🏢 商業機密", DANGER, "公司名稱、專案代號\n合作夥伴名單\n內部系統 URL"),
    ]:
        cx = ML + (cw + Inches(0.17)) * [("👤 個人資料", DANGER, ""),
             ("🔑 機密憑證", DANGER, ""), ("🏢 商業機密", DANGER, "")].index((title, bc, "")) # not ideal
        break

    cx = ML
    for title, bc, body in [
        ("👤 個人資料", DANGER, "姓名、Email\n電話、身分證\n客戶資料、合約對象"),
        ("🔑 機密憑證", DANGER, "API Key / Token\n資料庫密碼\nPrivate Key 路徑"),
        ("🏢 商業機密", DANGER, "公司名稱、專案代號\n合作夥伴名單\n內部系統 URL"),
    ]:
        card(slide, cx, cy, cw, ch, title=title, body=body,
             border=bc, title_color=bc, body_size=Pt(14))
        cx += cw + Inches(0.17)

    card(slide, ML, cy + ch + Inches(0.3), CW, Inches(1.1),
         body="⚠️  AI 服務的訓練政策、資料保留期限、第三方共享條款，都是你無法完全控制的。\n送出去的資料，就不再只屬於你。",
         border=DANGER, body_color=DANGER, body_size=Pt(15),
         fill=RGBColor(0x1c, 0x08, 0x08))


def s03_scenario(prs):
    slide = blank(prs)
    heading(slide, "真實情境：工程師請 AI 審查程式碼")

    before = ("# config.py\n"
              "DB_HOST     = \"prod-db.acme-corp.com\"\n"
              "DB_USER     = \"admin\"\n"
              "DB_PASSWORD = \"P@ssw0rd!2024\"\n"
              "API_KEY     = \"sk-live-abc123xyz789...\"\n\n"
              "OWNER_EMAIL = \"jason3e7@gmail.com\"\n"
              "PROJECT     = \"Project Phoenix\"")

    after = ("# config.py\n"
             "DB_HOST     = \"prod-db.[[KW_A1B2C3D4]].com\"\n"
             "DB_USER     = \"admin\"\n"
             "DB_PASSWORD = \"[[KW_E5F6A7B8]]\"\n"
             "API_KEY     = \"[[KW_C9D0E1F2]]\"\n\n"
             "OWNER_EMAIL = \"[[KW_33445566]]\"\n"
             "PROJECT     = \"[[KW_77889900]]\"")

    ch = Inches(3.8)
    code_block(slide, ML, BODY_TOP, HW, ch, before, "❌  直接貼上（原始）", label_color=DANGER)
    code_block(slide, R,  BODY_TOP, HW, ch, after,  "✅  用 kw-filter 處理後", label_color=ACCENT)

    txtbox(slide, ML, BODY_TOP + ch + Inches(0.2), CW, Inches(0.45),
           "AI 仍然能理解程式碼結構、提供重構建議——但不會看到任何真實的機密值。",
           size=Pt(14), color=MUTED)


def s04_commands(prs):
    slide = blank(prs)
    heading(slide, "運作原理：6 個獨立指令")

    headers = ["指令", "功能說明", "主要選項"]
    rows = [
        ["search",   "掃描目錄，列出關鍵字出現的檔案、行號、欄位（不分大小寫）", "-o FILE JSON 輸出  ·  --binary"],
        ["clear",    "將關鍵字替換為空字串或自訂文字，不需還原的場景",             "--replacement TEXT  ·  --dry-run  ·  --backup"],
        ["replace",  "換成 [[KW_XXXX]] token，同一關鍵字共用 token，輸出 mapping.json",  "-m FILE  ·  --dry-run  ·  --backup"],
        ["restore",  "讀取 mapping.json，將 token 換回原始值",                    "-m FILE  ·  --dry-run  ·  --backup"],
        ["cleanlog", "整行刪除含關鍵字的行，不留殘留結構",                        "--dry-run 預覽  ·  --stats 統計"],
        ["remap",    "依對照表替換值（IP / hostname / email → 假值），支援 binary 檔",  "--remap FILE  ·  --dry-run  ·  --backup"],
    ]
    table(slide, ML, BODY_TOP, CW, [1.2, 3.3, 2.7],
          headers, rows, row_h=Inches(0.48))

    flow_items = [
        ("📄 原始檔案", False, WHITE), ("→", False, MUTED),
        ("replace",    True,  ACCENT), ("→", False, MUTED),
        ("🤖 AI",      False, WHITE),  ("→", False, MUTED),
        ("AI 回應",    False, WHITE),  ("→", False, MUTED),
        ("restore",    True,  ACCENT), ("→", False, MUTED),
        ("✅ 還原",    False, WHITE),
    ]
    fy = BODY_TOP + Inches(0.48 * 7) + Inches(0.3)
    flow_row(slide, flow_items, ML + Inches(0.5), fy)


def s05_case1(prs):
    slide = blank(prs)
    heading(slide, "案例一：法律文件 AI 摘要")

    ty = BODY_TOP
    for label, color, bold, h in [
        ("情境",   ACCENT, True,  Inches(0.32)),
        ("律師事務所每週需要請 AI 摘要 30–50 份合約，合約內含客戶姓名、統編、金額、條款細節。",
                   WHITE,  False, Inches(0.5)),
        ("問題",   ACCENT, True,  Inches(0.32)),
        ("合約是客戶機密，直接貼上 ChatGPT 違反保密義務，且違反 GDPR、個資法等資料保護法規。",
                   WHITE,  False, Inches(0.5)),
        ("解法",   ACCENT, True,  Inches(0.32)),
        ("維護 keywords.txt，執行 replace 再送 AI；拿回摘要後執行 restore 還原。",
                   WHITE,  False, Inches(0.45)),
    ]:
        txtbox(slide, ML, ty, HW, h, label, size=Pt(15) if bold else Pt(14),
               color=color, bold=bold)
        ty += h + Inches(0.06)

    before = ("甲方：台灣科技股份有限公司\n"
              "代表人：陳建宏\n"
              "統一編號：12345678\n"
              "合約金額：新台幣 500 萬元\n"
              "生效日期：2024-03-01")
    after  = ("甲方：[[KW_A1B2C3D4]]\n"
              "代表人：[[KW_E5F6A7B8]]\n"
              "統一編號：[[KW_C9D0E1F2]]\n"
              "合約金額：新台幣 [[KW_33445566]] 元\n"
              "生效日期：2024-03-01")
    code_block(slide, R, BODY_TOP,              HW, Inches(1.95), before, "送出前（原始）",    label_color=DANGER)
    code_block(slide, R, BODY_TOP + Inches(2.1),HW, Inches(1.95), after,  "送出後（token 化）",label_color=ACCENT)

    txtbox(slide, ML, Inches(5.55), CW, Inches(0.38),
           "AI 仍能摘要條款結構、義務事項、違約責任——機密數值完全遮蔽。",
           size=Pt(13), color=MUTED)


def s06_case2(prs):
    slide = blank(prs)
    heading(slide, "案例二：程式碼庫 AI 重構")

    ty = BODY_TOP
    txtbox(slide, ML, ty, HW, Inches(0.32), "情境", size=Pt(15), color=ACCENT, bold=True)
    ty += Inches(0.35)
    txtbox(slide, ML, ty, HW, Inches(0.44),
           "後端團隊想用 AI 重構 200+ 個 Python 檔，但 repo 中散落著：",
           size=Pt(14), color=WHITE)
    ty += Inches(0.47)

    for b in ["硬編碼的 API Key", "內部服務 domain（api.internal.acme.com）",
              "資料庫連線字串", "員工姓名寫在 TODO 裡"]:
        txtbox(slide, ML + Inches(0.2), ty, HW, Inches(0.33),
               "•  " + b, size=Pt(13.5), color=WHITE)
        ty += Inches(0.33)

    ty += Inches(0.12)
    txtbox(slide, ML, ty, HW, Inches(0.32), "執行方式", size=Pt(15), color=ACCENT, bold=True)
    ty += Inches(0.35)

    code = ("# 1. 確認覆蓋範圍\n"
            "python3 kw_tools.py search  -k kw.txt -t ./src -r\n"
            "# 2. 替換為 token\n"
            "python3 kw_tools.py replace -k kw.txt -t ./src -r -m map.json\n"
            "# 3. AI 回應後還原\n"
            "python3 kw_tools.py restore -m map.json -t ./out -r")
    code_block(slide, ML, ty, HW, Inches(1.75), code)

    # Right
    txtbox(slide, R, BODY_TOP, HW, Inches(0.32), "效果", size=Pt(15), color=ACCENT, bold=True)
    ey = BODY_TOP + Inches(0.38)
    for i, body in enumerate([
        "✅  AI 看到完整的邏輯結構、函式命名、模組關係",
        "✅  所有機密字串以 token 形式保留，不影響 AI 理解程式流程",
        "✅  還原後，重構結果直接可用，不需要手動把 API Key 貼回去",
        "📊  實測：100 個檔案 / 10 萬行  →  處理時間  0.34 秒",
    ]):
        bc = BLUE if i == 3 else ACCENT
        card(slide, R, ey, HW, Inches(1.05),
             body=body, border=bc, body_size=Pt(14))
        ey += Inches(1.13)


def s07_case3(prs):
    slide = blank(prs)
    heading(slide, "案例三：客服日誌 AI 分析")

    ty = BODY_TOP
    txtbox(slide, ML, ty, HW, Inches(0.32), "情境", size=Pt(15), color=ACCENT, bold=True)
    ty += Inches(0.35)
    txtbox(slide, ML, ty, HW, Inches(0.5),
           "電商公司每月有 5 萬筆客服對話紀錄，想用 AI 分析常見問題與情緒趨勢，但日誌中含有：",
           size=Pt(14), color=WHITE)
    ty += Inches(0.54)

    for b in ["用戶姓名、Email", "訂單編號（可反查個資）", "客服人員姓名", "地址片段"]:
        txtbox(slide, ML + Inches(0.2), ty, HW, Inches(0.33),
               "•  " + b, size=Pt(13.5), color=WHITE)
        ty += Inches(0.33)

    ty += Inches(0.15)
    txtbox(slide, ML, ty, HW, Inches(0.32), "Workflow", size=Pt(15), color=ACCENT, bold=True)
    ty += Inches(0.36)

    for step in ["① search  — 確認關鍵字覆蓋率",
                 "② replace — 批次處理 5 萬筆日誌",
                 "③ 送 AI 做分類、情緒分析、趨勢報告",
                 "④ restore — 報告中如有引用個案，還原姓名"]:
        shape = rect(slide, ML, ty, HW, Inches(0.44))
        tf = shape.text_frame
        _set_txbody_margins(tf._txBody, l=Pt(10), t=Pt(5), anchor='ctr')
        p = tf.paragraphs[0]
        run = p.add_run(); run.text = step
        run.font.size = Pt(13.5); run.font.color.rgb = WHITE; run.font.name = "Noto Sans TC"
        ty += Inches(0.48)

    # Right
    txtbox(slide, R, BODY_TOP, HW, Inches(0.32), "關鍵字清單範例",
           size=Pt(15), color=ACCENT, bold=True)
    kw = ("# customer_kw.txt\n# 常見姓名\n王小明 / 陳美麗 / 李大偉\n"
          "# 客服人員\n服務專員 Andy\n服務專員 Sarah\n# 訂單前綴\nORD-2024-")
    code_block(slide, R, BODY_TOP + Inches(0.38), HW, Inches(2.3), kw)

    card(slide, R, BODY_TOP + Inches(2.78), HW, Inches(0.9),
         body="符合個人資料保護法第 6 條「去識別化」處理要求，可作為合規佐證。",
         border=ACCENT, body_size=Pt(14))


def s08_cleanlog(prs):
    slide = blank(prs)
    heading(slide, "案例四：系統日誌清理（cleanlog 新功能）")

    ty = BODY_TOP
    txtbox(slide, ML, ty, HW, Inches(0.32), "情境", size=Pt(15), color=ACCENT, bold=True)
    ty += Inches(0.35)
    txtbox(slide, ML, ty, HW, Inches(0.5),
           "DevOps 工程師需要把 production log 送給 AI 找出異常模式，但日誌中夾雜：",
           size=Pt(14), color=WHITE)
    ty += Inches(0.54)

    for b in ["用戶 Email / 帳號 ID", "IP 位址（內部拓樸資訊）",
              "含機密值的 debug 行（token=xxx）", "認證錯誤行（含嘗試的密碼明文）"]:
        txtbox(slide, ML + Inches(0.2), ty, HW, Inches(0.33),
               "•  " + b, size=Pt(13.5), color=WHITE)
        ty += Inches(0.33)

    ty += Inches(0.15)
    txtbox(slide, ML, ty, HW, Inches(0.42),
           "這類日誌「部分遮蔽」不夠——只要含有機密，整行都需要移除。",
           size=Pt(14), color=WARN)
    ty += Inches(0.45)

    card(slide, ML, ty, HW, Inches(1.05),
         title="cleanlog 與 clear 的差異",
         body="clear：token=secret → token=（殘留結構）\ncleanlog：整行刪除，不留任何痕跡",
         border=PURPLE, title_color=PURPLE, body_size=Pt(13.5))

    # Right
    raw = ("10:00:01 INFO  Server started\n"
           "10:00:05 INFO  Login: alice@corp.com\n"
           "10:01:22 ERROR Auth fail: secret_token invalid\n"
           "10:01:23 INFO  Retry connection...\n"
           "10:02:00 ERROR Timeout on 192.168.1.100\n"
           "10:02:01 INFO  Connection restored\n"
           "10:03:10 DEBUG token=sk-live-abc123 cached")
    clean = ("10:00:01 INFO  Server started\n"
             "10:01:23 INFO  Retry connection...\n"
             "10:02:01 INFO  Connection restored")
    cmd = ("# 預覽模式（不修改檔案）\n"
           "python3 kw_tools.py cleanlog -k kw.txt -t ./logs -r --dry-run\n\n"
           "# 實際執行（含統計）\n"
           "python3 kw_tools.py cleanlog -k kw.txt -t ./logs -r --stats")

    code_block(slide, R, BODY_TOP,               HW, Inches(2.1),  raw,   "原始 log（含機密）",  label_color=DANGER)
    code_block(slide, R, BODY_TOP + Inches(2.2),  HW, Inches(1.05), clean, "cleanlog 處理後",     label_color=ACCENT)
    code_block(slide, R, BODY_TOP + Inches(3.35), HW, Inches(1.35), cmd)


def s08b_playwright(prs):
    slide = blank(prs)
    heading(slide, "案例五：登入後 HTML → 過濾 PII → AI 生成 Playwright 測試")

    ty = BODY_TOP
    txtbox(slide, ML, ty, HW, Inches(0.32), "情境", size=Pt(15), color=ACCENT, bold=True)
    ty += Inches(0.36)
    txtbox(slide, ML, ty, HW, Inches(0.52),
           "QA 工程師登入內部 CRM，將儀表板 HTML 存下來，準備讓 AI 產生 Playwright E2E 測試。"
           "但 HTML 中混有大量真實機敏資料：",
           size=Pt(13.5), color=WHITE)
    ty += Inches(0.58)

    for b in ["客戶 PII：姓名、Email、電話（wang.daming@gmail.com / 0912-345-678）",
              "登入帳號（alice.wu@acme-corp.internal）",
              "Session JWT（eyJhbGci…acme2024）",
              "API Key（sk-prod-4f8e2a7b9c1d3e5f）",
              "內部主機名稱（acme-corp.internal）"]:
        txtbox(slide, ML + Inches(0.2), ty, HW, Inches(0.3),
               "•  " + b, size=Pt(12.5), color=WHITE)
        ty += Inches(0.3)

    ty += Inches(0.12)
    txtbox(slide, ML, ty, HW, Inches(0.32), "Workflow", size=Pt(15), color=ACCENT, bold=True)
    ty += Inches(0.35)

    for step in ["① replace  — HTML 中所有機敏值換成 token",
                 "② 將 token 化的 HTML 送給 AI",
                 "③ AI 產出含 token 的 Playwright TypeScript",
                 "④ restore  — token 換回真實值，腳本可直接跑"]:
        shape = rect(slide, ML, ty, HW, Inches(0.40))
        tf = shape.text_frame
        _set_txbody_margins(tf._txBody, l=Pt(10), t=Pt(4), anchor='ctr')
        p = tf.paragraphs[0]
        run = p.add_run(); run.text = step
        run.font.size = Pt(12.5); run.font.color.rgb = WHITE; run.font.name = "Noto Sans TC"
        ty += Inches(0.44)

    # Right column: 3 × 2.05" code blocks, each 6 lines, total fits within 7.5"
    # label=0.3" + rect=1.75", gap between blocks=0.1"
    # ends at BODY_TOP + 3*2.05 + 2*0.1 = 1.18+6.35 = 7.53 → use 2.0" each
    token_html = ('<span id="current-user">[[KW_A1B2C3D4]]</span>\n'
                  '<td>[[KW_E5F6A7B8]]</td>\n'
                  '<td>[[KW_C9D0E1F2]]</td>\n'
                  'const SESSION = "[[KW_11223344]]";\n'
                  'const BASE_URL =\n'
                  '  "https://api.[[KW_55667788]]:8443/v2";')

    ai_code =    ("// AI 產出（含 token）\n"
                  "await page.goto(\n"
                  "  'https://crm.[[KW_55667788]]/dashboard');\n"
                  "await expect(page.locator('#current-user'))\n"
                  "  .toHaveText('[[KW_A1B2C3D4]]');\n"
                  "// row C001 → [[KW_E5F6A7B8]]")

    final_code = ("// restore 後（可直接執行）\n"
                  "await page.goto(\n"
                  "  'https://crm.acme-corp.internal/dashboard');\n"
                  "await expect(page.locator('#current-user'))\n"
                  "  .toHaveText('alice.wu@acme-corp.internal');\n"
                  "// row C001 → wang.daming@gmail.com")

    BH = Inches(2.0)   # block height (label 0.3 + rect 1.7)
    BG = Inches(0.1)   # gap between blocks
    code_block(slide, R, BODY_TOP,              HW, BH, token_html, "① replace — HTML（節錄）", label_color=WARN)
    code_block(slide, R, BODY_TOP + BH + BG,    HW, BH, ai_code,    "③ AI 回傳（token）",       label_color=MUTED)
    code_block(slide, R, BODY_TOP + 2*(BH+BG),  HW, BH, final_code, "④ restore 後",             label_color=ACCENT)


def s09_tests(prs):
    slide = blank(prs)
    heading(slide, "驗證一：自動化測試覆蓋")

    headers = ["測試模組", "數量", "涵蓋面向"]
    rows = [
        ["test_utils.py",              "22",  "工具函式、Regex 編譯、Binary search"],
        ["test_search.py",             "16",  "搜尋、多檔案、遞迴、Unicode、JSON 輸出"],
        ["test_clear.py",              "14",  "清空、自訂替換文字、backup、多檔"],
        ["test_replace.py",            "14",  "Token 格式、穩定性、最長關鍵字優先"],
        ["test_restore.py",            "13",  "還原、Roundtrip、錯誤處理"],
        ["test_cleanlog.py",           "22",  "整行刪除、dry-run、stats、Unicode"],
        ["test_remap.py",              "20",  "Binary 模式替換、dry-run、Unicode"],
        ["test_dry_run_ignore_case.py","18",  "dry-run（各指令）、大小寫不分"],
        ["test_playwright_scenario.py", "11", "HTML PII 過濾 → AI → restore 端對端"],
        ["test_integration.py",        "14",  "CLI subprocess 端對端測試"],
        ["合計",                       "164", ""],
    ]
    table(slide, ML, BODY_TOP, HW, [3.0, 0.65, 3.0],
          headers, rows, row_h=Inches(0.48), cell_size=Pt(11))

    pytest_out = ("$ python3 -m pytest tests/ -v\n\n"
                  "tests/test_utils.py               ✓ 22\n"
                  "tests/test_search.py              ✓ 16\n"
                  "tests/test_clear.py               ✓ 14\n"
                  "tests/test_replace.py             ✓ 14\n"
                  "tests/test_restore.py             ✓ 13\n"
                  "tests/test_cleanlog.py            ✓ 22\n"
                  "tests/test_remap.py               ✓ 20\n"
                  "tests/test_dry_run_ignore_case.py ✓ 18\n"
                  "tests/test_playwright_scenario.py ✓ 11\n"
                  "tests/test_integration.py         ✓ 14\n\n"
                  "======= 164 passed in 2.1s =======")
    code_block(slide, R + Inches(0.15), BODY_TOP, HW - Inches(0.15), Inches(3.45), pytest_out)

    card(slide, R + Inches(0.15), BODY_TOP + Inches(3.6), HW - Inches(0.15), Inches(1.6),
         body=("Roundtrip（test_restore.py）：replace → restore 後逐字元比對，確保 100% 還原。\n\n"
               "Playwright 情境（test_playwright_scenario.py）：登入後 HTML → PII 過濾 → AI 生成 → restore，驗證無 token 殘留。"),
         border=ACCENT, body_size=Pt(12.5))


def s10_perf(prs):
    slide = blank(prs)
    heading(slide, "驗證二：效能實測")

    sw = Inches(3.85); sh = Inches(1.5); sy = BODY_TOP
    for num, lbl, sx in [("0.34s", "100 檔 / 10 萬行", ML),
                          ("292,000", "行/秒（replace）", ML + sw + Inches(0.24)),
                          ("1", "Regex 編譯次數",   ML + (sw + Inches(0.24)) * 2)]:
        stat_box(slide, sx, sy, sw, sh, num, lbl)

    card_y = BODY_TOP + sh + Inches(0.35)
    card_h = Inches(2.65)
    card(slide, ML, card_y, HW, card_h,
         title="為什麼快？",
         body=("• 所有關鍵字合併成單一 Regex，只掃描一次\n"
               "• 關鍵字依長度排序，最長優先，避免重複回溯\n"
               "• Binary search（bisect）做 O(log n) 關鍵字存在確認\n"
               "• 逐行 iterator，不一次載入整個大檔"),
         body_size=Pt(14))

    card(slide, R, card_y, HW, card_h,
         title="測試環境",
         body=("• Python 3.10 · Ubuntu 20.04\n"
               "• 4 個關鍵字，每行平均含 2 個命中\n"
               "• 100 個檔案 × 1,000 行 = 100,000 行\n"
               "• 無 SSD 加速（標準 VM 磁碟）\n\n"
               "實際企業文件通常每行命中率更低，速度會更快。"),
         body_size=Pt(14))


def s11_correctness(prs):
    slide = blank(prs)
    heading(slide, "驗證三：正確性保證")

    cw = Inches(3.9); cy = BODY_TOP; ch = Inches(3.0); gap = Inches(0.22)
    cx = ML
    for title, desc, code_txt in [
        ("最長優先",
         "關鍵字 John 和 John Doe 同時存在時，John Doe 優先被匹配，不會拆成兩個 token。",
         "Input:  John Doe was here\nOutput: [[KW_AAAA]] was here\n✓ 只有 1 個 token（非 2 個）"),
        ("穩定 Token",
         "同一關鍵字在所有檔案中對應同一個 token，還原後保持語意一致性。",
         "file_a.txt: [[KW_AAAA]]\nfile_b.txt: [[KW_AAAA]]\n→ restore → 兩者都還原為 \"John Doe\""),
        ("Unicode 完整支援",
         "繁體中文、日文、特殊符號均正確處理，不會造成字元截斷或亂碼。",
         "Input:  作者：張三\nOutput: 作者：[[KW_BBBB]]\nrestore → 作者：張三 ✓"),
    ]:
        shape = rect(slide, cx, cy, cw, ch, fill_color=BG2, border_color=ACCENT)
        tf = shape.text_frame; tf.word_wrap = True
        _set_txbody_margins(tf._txBody, l=Pt(12), t=Pt(11))
        # title
        p = tf.paragraphs[0]; r = p.add_run()
        r.text = title; r.font.size = Pt(16); r.font.bold = True
        r.font.color.rgb = ACCENT; r.font.name = "Noto Sans TC"
        # desc
        p2 = tf.add_paragraph(); p2.space_before = Pt(6)
        r2 = p2.add_run(); r2.text = desc
        r2.font.size = Pt(13.5); r2.font.color.rgb = WHITE; r2.font.name = "Noto Sans TC"
        # code
        p3 = tf.add_paragraph(); p3.space_before = Pt(10)
        r3 = p3.add_run(); r3.text = code_txt
        r3.font.size = Pt(11.5); r3.font.color.rgb = CODE_FG; r3.font.name = "Courier New"
        cx += cw + gap

    card(slide, ML, cy + ch + Inches(0.28), CW, Inches(0.8),
         body="🔁  Roundtrip 保證：replace → restore 後，所有測試案例均通過逐字元比對（含多檔案、Unicode、多次出現同一關鍵字）。",
         border=ACCENT, body_color=ACCENT, body_size=Pt(15),
         fill=RGBColor(0x0f, 0x1f, 0x12))


def s12_comparison(prs):
    slide = blank(prs)
    heading(slide, "與現有方案比較")

    headers = ["方案", "安裝複雜度", "可還原", "自訂關鍵字", "CLI 可用", "適用場景"]
    rows = [
        ["kw-filter（本專案）",  "零依賴",           "✅ mapping.json", "✅ 完全自訂",  "✅ 5 個指令",   "文字檔、程式碼、日誌"],
        ["Microsoft Presidio",  "重（spaCy model）", "✅ 支援",         "⚠️ 需設定",   "⚠️ Python API", "企業 PII 偵測"],
        ["手動 sed / replace",  "零",               "❌ 不可還原",     "✅",          "✅",            "簡單一次性需求"],
        ["LangChain Anonymizer","重（langchain）",   "✅",             "⚠️",         "❌ 需整合",     "LangChain pipeline"],
        ["人工遮蔽",            "零",               "❌",             "✅",          "❌",            "少量、一次性"],
    ]
    table(slide, ML, BODY_TOP, CW, [2.2, 1.75, 1.75, 1.55, 1.65, 2.4],
          headers, rows, row_h=Inches(0.55))

    note_y = BODY_TOP + Inches(0.55 * 6) + Inches(0.35)
    txtbox(slide, ML, note_y, CW, Inches(0.42),
           "kw-filter 定位：輕量、可還原、自訂度高，適合開發者和小型團隊快速整合到現有 workflow。",
           size=Pt(14), color=MUTED)


def s13_why(prs):
    slide = blank(prs)
    heading(slide, "為什麼選擇 kw-filter？")

    items = [
        ("🔒 資料不離開你的掌控",  "所有處理都在本地執行，不呼叫任何外部服務。mapping.json 只存在你的機器上。"),
        ("⚡ 5 分鐘內可上線",       "只需要 Python 3.8+，git clone 後直接執行。不需要安裝 model、不需要設定帳號。"),
        ("🔁 無損還原",             "token 化後 AI 的回應同樣包含 token，執行 restore 後完整還原，無需手動替換。"),
        ("📋 可稽核、可合規",       "mapping.json 是明確的處理紀錄，可作為 GDPR / 個資法去識別化處理的文件佐證。"),
        ("🧩 適合任何 AI 服務",     "不綁定特定 AI。ChatGPT、Claude、Gemini、本地 LLM 都適用——就是純文字替換。"),
        ("🧪 160+ 個測試保證品質",  "涵蓋 edge case：Unicode、多檔案、最長優先、Roundtrip 還原、Playwright 情境 E2E。"),
    ]
    ch = Inches(1.62)
    cx_l, cx_r = ML, R
    cy = BODY_TOP
    for i, (title, body) in enumerate(items):
        cx = cx_l if i % 2 == 0 else cx_r
        if i % 2 == 0 and i > 0:
            cy += ch + Inches(0.14)
        card(slide, cx, cy, HW, ch, title=title, body=body,
             border=ACCENT, body_size=Pt(14))


def s14_quickstart(prs):
    slide = blank(prs)
    heading(slide, "立刻開始使用")

    # Left column
    install = "git clone https://github.com/jason3e7/kw-filter.git\ncd kw-filter"
    kw_file = "# keywords.txt\nJohn Doe\nacme-corp.com\nsk-live-abc123"
    run_cmd = ("# 遞迴 + 大小寫不分均為預設，無需額外參數\n"
               "python3 kw_tools.py search   -k kw.txt -t ./docs\n"
               "python3 kw_tools.py replace  -k kw.txt -t ./docs\n"
               "python3 kw_tools.py restore  -t ./ai_output\n"
               "python3 kw_tools.py cleanlog -k kw.txt -t ./logs --stats\n"
               "python3 kw_tools.py remap    --remap remap.txt -t ./logs")

    code_block(slide, ML, BODY_TOP,                HW, Inches(1.05), install, "安裝")
    code_block(slide, ML, BODY_TOP + Inches(1.17), HW, Inches(1.55), kw_file, "建立關鍵字清單")
    code_block(slide, ML, BODY_TOP + Inches(2.84), HW, Inches(2.05), run_cmd, "執行")

    # Right column — online tool card + test + summary
    # Online tool card (highlight box)
    online_box = rect(slide, R, BODY_TOP, HW, Inches(1.25),
                      fill_color=RGBColor(0x0a, 0x24, 0x18), border_color=ACCENT, border_pt=2)
    tf_o = online_box.text_frame; tf_o.word_wrap = True
    _set_txbody_margins(tf_o._txBody, l=Pt(13), t=Pt(10))
    p_o = tf_o.paragraphs[0]; p_o.alignment = PP_ALIGN.LEFT
    ro1 = p_o.add_run()
    ro1.text = "🌐  線上互動工具（無需安裝）"
    ro1.font.size = Pt(14); ro1.font.bold = True
    ro1.font.color.rgb = ACCENT; ro1.font.name = "Noto Sans TC"
    p_o2 = tf_o.add_paragraph(); p_o2.space_before = Pt(5); p_o2.alignment = PP_ALIGN.LEFT
    ro2 = p_o2.add_run()
    ro2.text = "jason3e7.github.io/kw-filter/kw_tools.html"
    ro2.font.size = Pt(12.5); ro2.font.color.rgb = BLUE; ro2.font.name = "Courier New"
    p_o3 = tf_o.add_paragraph(); p_o3.space_before = Pt(4); p_o3.alignment = PP_ALIGN.LEFT
    ro3 = p_o3.add_run()
    ro3.text = "載入檔案 → 選擇操作 → 下載結果，不需 Python 環境"
    ro3.font.size = Pt(12); ro3.font.color.rgb = MUTED; ro3.font.name = "Noto Sans TC"

    # Test block
    test_code = "pip install pytest\npython3 -m pytest tests/ -v"
    code_block(slide, R, BODY_TOP + Inches(1.38), HW, Inches(1.05), test_code, "執行測試")

    # Summary card
    shape = rect(slide, R, BODY_TOP + Inches(2.55), HW, Inches(1.84),
                 fill_color=RGBColor(0x0f, 0x1f, 0x12), border_color=ACCENT)
    tf = shape.text_frame; tf.word_wrap = True
    _set_txbody_margins(tf._txBody, l=Pt(13), t=Pt(10))

    summary_title = tf.paragraphs[0]
    summary_title.alignment = PP_ALIGN.LEFT
    rt = summary_title.add_run()
    rt.text = "重點回顧"; rt.font.size = Pt(14); rt.font.bold = True
    rt.font.color.rgb = ACCENT; rt.font.name = "Noto Sans TC"

    for item in ["6 個 CLI 指令 + GitHub Pages 線上工具",
                 "遞迴 & 大小寫不分均為預設行為",
                 "replace → restore  100% 還原",
                 "160+ 個測試，零第三方依賴",
                 "Playwright 情境：爬蟲 → AI → 可執行腳本"]:
        p = tf.add_paragraph(); p.space_before = Pt(4)
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run(); run.text = "▸  " + item
        run.font.size = Pt(12.5); run.font.color.rgb = WHITE; run.font.name = "Noto Sans TC"

    txtbox(slide, ML, Inches(6.55), CW, Inches(0.3),
           "github.com/jason3e7/kw-filter",
           size=Pt(11.5), color=MUTED, align=PP_ALIGN.RIGHT)


def s14b_mcp(prs):
    slide = blank(prs)
    heading(slide, "MCP 整合：Claude Code × kw-filter Server")

    # ── Left column ───────────────────────────────────────────────────────────
    # Architecture: 6 lines → fits in 1.80" total (incl. label)
    arch = (
        "Claude Code ──stdio──► client.py\n"
        "                         │  HTTP :8000\n"
        "                     server.py :8000\n"
        "                ┌────────┴────────┐\n"
        "           storage/          restored/\n"
        "         (tokenised)         (還原後)"
    )
    code_block(slide, ML, BODY_TOP, HW, Inches(1.80), arch, "架構", label_color=ACCENT)

    # 3 MCP tool cards — height 0.65" each, step 0.72"
    tools_y = BODY_TOP + Inches(1.92)
    card(slide, ML, tools_y, HW, Inches(0.65),
         title="list_files()",
         body="列出所有已上傳（自動 tokenised）的檔案",
         body_size=Pt(12.5))
    card(slide, ML, tools_y + Inches(0.72), HW, Inches(0.65),
         title="get_files(file_id)",
         body="取得 tokenised 內容 — 安全送入任何 AI",
         body_size=Pt(12.5))
    card(slide, ML, tools_y + Inches(1.44), HW, Inches(0.65),
         title="upload_files(name, content)",
         body="AI 產出含 token 的內容 → 自動還原並儲存",
         body_size=Pt(12.5))

    # Claude config snippet — 6 lines
    cfg = (
        '# ~/.claude.json\n'
        '"kw-filter": {\n'
        '  "command": "python3",\n'
        '  "args": [".../mcp/client.py"],\n'
        '  "env": {"KW_SERVER_URL": "http://localhost:8000"}\n'
        '}'
    )
    code_block(slide, ML, tools_y + Inches(2.21), HW, Inches(1.80), cfg, "Claude Code 設定")

    # ── Right column ──────────────────────────────────────────────────────────
    # Endpoints table — 9 data rows at 0.35" each
    tbl_rows = [
        ["GET  /",                   "Web UI — 瀏覽器管理介面"],
        ["POST /files  /files/text", "上傳並自動 replace，回傳 file_id"],
        ["GET  /files",              "列出已 tokenised 的檔案"],
        ["GET  /files/{id}",         "取得 tokenised 內容"],
        ["DELETE /files/{id}",       "刪除 tokenised 檔案"],
        ["POST /restore",            "還原 token → 原始值，存入 restored/"],
        ["GET  /restored",           "列出所有還原後的檔案"],
        ["GET|DELETE /restored/{n}", "下載／刪除還原後的檔案"],
        ["GET|PUT /keywords",        "查看／更新 keywords.txt"],
    ]
    table(slide, R, BODY_TOP, HW,
          col_fracs=[1.8, 2.2],
          headers=["Endpoint", "說明"],
          rows=tbl_rows,
          row_h=Inches(0.35),
          cell_size=Pt(10.5))

    # IP blacklist: table height = 10 rows * 0.35 = 3.50"
    bl_y = BODY_TOP + Inches(3.62)
    blacklist_code = (
        "# mcp/ip_blacklist.txt\n"
        "# 每行一個 IP，# 開頭為註解\n"
        "# 空白檔案 = 不限制任何 IP\n"
        "203.0.113.10\n"
        "198.51.100.42"
    )
    code_block(slide, R, bl_y, HW, Inches(1.38), blacklist_code, "IP 黑名單設定")

    # Restricted endpoints note
    note_y = bl_y + Inches(1.49)
    card(slide, R, note_y, HW, Inches(0.62),
         body="黑名單 IP 無法存取：/docs · /keywords · /restored\n檔案即時生效，無需重啟 server",
         body_size=Pt(12), fill=RGBColor(0x1a, 0x0d, 0x0d),
         border=DANGER)

    txtbox(slide, ML, Inches(6.90), CW, Inches(0.3),
           "github.com/jason3e7/kw-filter",
           size=Pt(11.5), color=MUTED, align=PP_ALIGN.RIGHT)


def s14c_webui(prs):
    slide = blank(prs)
    heading(slide, "Web UI：瀏覽器管理介面（GET /）")

    CARD_BG   = RGBColor(0x12, 0x20, 0x18)
    CARD_BD   = RGBColor(0x1e, 0x3a, 0x28)
    cw = Inches(2.77)   # 2 cards per column, gap 0.15"
    ch = Inches(1.52)
    gap = Inches(0.15)

    features = [
        ("📤", "上傳檔案",
         "拖曳或選擇本機檔案（支援多選）\n"
         "或直接貼上文字內容\n"
         "上傳後自動執行 replace"),
        ("📁", "Tokenised 檔案",
         "列出所有已上傳的檔案\n"
         "顯示 token 替換數量\n"
         "可下載或刪除"),
        ("✅", "已還原檔案",
         "列出所有 restore 完成的檔案\n"
         "可下載原始內容或刪除\n"
         "（不出現在 list_files）"),
        ("🔄", "還原表單",
         "貼上含 [[KW_...]] 的 AI 輸出\n"
         "一鍵 restore 並存檔\n"
         "無需 CLI 或 MCP"),
        ("🔑", "Keywords 管理",
         "在瀏覽器中查看和編輯\n"
         "keywords.txt 內容\n"
         "儲存後立即生效"),
        ("⚡", "純 HTML / JS",
         "無框架、無需額外安裝\n"
         "由 server.py 直接服務\n"
         "深色主題，響應式佈局"),
    ]

    col_x = [ML, ML + cw + gap]
    row_y = [BODY_TOP, BODY_TOP + ch + gap, BODY_TOP + 2 * (ch + gap)]

    for i, (icon, title, body) in enumerate(features):
        cx = col_x[i % 2]
        cy = row_y[i // 2]

        shape = rect(slide, cx, cy, cw, ch, fill_color=CARD_BG, border_color=CARD_BD, border_pt=1.5)
        tf = shape.text_frame; tf.word_wrap = True
        _set_txbody_margins(tf._txBody, l=Pt(12), r=Pt(10), t=Pt(10), b=Pt(8))

        p0 = tf.paragraphs[0]
        ri = p0.add_run(); ri.text = icon + "  "
        ri.font.size = Pt(16); ri.font.name = "Segoe UI Emoji"
        rt = p0.add_run(); rt.text = title
        rt.font.size = Pt(14); rt.font.bold = True
        rt.font.color.rgb = ACCENT; rt.font.name = "Noto Sans TC"

        p1 = tf.add_paragraph(); p1.space_before = Pt(6)
        rb = p1.add_run(); rb.text = body
        rb.font.size = Pt(12); rb.font.color.rgb = WHITE; rb.font.name = "Noto Sans TC"

    # ── Right column: access + workflow ──────────────────────────────────────
    rx = ML + 2 * cw + 2 * gap + Inches(0.05)
    rw = W - rx - ML

    access_code = (
        "# 啟動 server 後直接開啟瀏覽器\n"
        "python mcp/server.py\n"
        "# → http://localhost:8000"
    )
    code_block(slide, rx, BODY_TOP, rw, Inches(1.10), access_code, "存取方式", label_color=ACCENT)

    # Workflow steps
    steps = [
        ("① 上傳",    "拖曳或貼上文字 → 自動 replace"),
        ("② 複製",    "下載或複製 tokenised 內容"),
        ("③ 送 AI",   "含 [[KW_...]] 的安全內容送入 AI"),
        ("④ 還原",    "貼上 AI 輸出 → 一鍵 restore"),
        ("⑤ 下載",    "從「已還原檔案」取得最終結果"),
    ]
    wy = BODY_TOP + Inches(1.22)
    txtbox(slide, rx, wy, rw, Inches(0.28), "典型工作流程",
           size=Pt(11), color=MUTED)
    wy += Inches(0.30)
    for step, desc in steps:
        shape = rect(slide, rx, wy, rw, Inches(0.56),
                     fill_color=RGBColor(0x0e, 0x1c, 0x12), border_color=BORDER)
        tf = shape.text_frame; tf.word_wrap = True
        _set_txbody_margins(tf._txBody, l=Pt(10), r=Pt(8), t=Pt(7), b=Pt(7))
        p = tf.paragraphs[0]
        rs = p.add_run(); rs.text = step + "  "
        rs.font.size = Pt(12); rs.font.bold = True
        rs.font.color.rgb = ACCENT; rs.font.name = "Noto Sans TC"
        rd = p.add_run(); rd.text = desc
        rd.font.size = Pt(12); rd.font.color.rgb = WHITE; rd.font.name = "Noto Sans TC"
        wy += Inches(0.62)

    txtbox(slide, ML, Inches(6.90), CW, Inches(0.3),
           "github.com/jason3e7/kw-filter",
           size=Pt(11.5), color=MUTED, align=PP_ALIGN.RIGHT)


def s15_bigquestion(prs):
    slide = blank(prs)
    fill = slide.background.fill; fill.solid()
    fill.fore_color.rgb = RGBColor(0x0d, 0x04, 0x1e)

    txtbox(slide, ML, Inches(0.38), CW, Inches(0.32),
           "大 哉 問",
           size=Pt(14), color=PURPLE, align=PP_ALIGN.CENTER)

    txtbox(slide, ML, Inches(0.72), CW, Inches(0.95),
           "AI 讓我們隨時隨地都能工作，\n但我們真的準備好「隨時隨地都在工作」了嗎？",
           size=Pt(26), color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    div = rect(slide, Inches(6.17), Inches(1.75), Inches(1.0), Pt(3),
               fill_color=PURPLE, border_color=None)
    div.line.fill.background()

    # 2 × 2 card grid
    CARD_BG     = RGBColor(0x18, 0x08, 0x2a)
    CARD_BORDER = RGBColor(0x5e, 0x28, 0x8e)
    CARD_TEXT   = RGBColor(0xcb, 0xd5, 0xe1)

    cw = HW; ch = Inches(1.95); gap_x = GAP; gap_y = Inches(0.2)
    row1_y = Inches(1.92); row2_y = row1_y + ch + gap_y

    cards = [
        # row 1
        (ML,  row1_y, "🌙 工作侵入生活",
         "AI 讓凌晨兩點整理報告變得「輕鬆」。\n"
         "輕鬆完成，不等於應該完成。\n"
         "工具降低了摩擦，卻也模糊了「下班」這條線。"),
        (R,   row1_y, "📱 手機 × AI ＝ 永遠在線",
         "手機讓你隨時隨地都能用 AI——\n"
         "通勤、吃飯、睡前都能繼續「生產」。\n"
         "感覺一直在工作，因為你確實一直在工作。"),
        # row 2
        (ML,  row2_y, "😰 額度焦慮：全新的壓力",
         "用完了 → 今天無法繼續，焦慮。\n"
         "沒用完 → 浪費配額，也焦慮。\n"
         "一種以前從未存在過的、工具獨有的心理負擔。"),
        (R,   row2_y, "🧠 決策疲勞轉移",
         "AI 接管了大量判斷，\n"
         "但「要不要用 AI 做這件事」本身就是新的負擔。\n"
         "我們節省了時間，卻多了另一種焦慮。"),
    ]

    for cx, cy, title, body in cards:
        shape = rect(slide, cx, cy, cw, ch,
                     fill_color=CARD_BG,
                     border_color=CARD_BORDER, border_pt=1.5)
        tf = shape.text_frame; tf.word_wrap = True
        _set_txbody_margins(tf._txBody, l=Pt(12), r=Pt(12), t=Pt(10), b=Pt(8))
        p = tf.paragraphs[0]; r = p.add_run()
        r.text = title; r.font.size = Pt(14.5); r.font.bold = True
        r.font.color.rgb = PURPLE; r.font.name = "Noto Sans TC"
        p2 = tf.add_paragraph(); p2.space_before = Pt(6)
        r2 = p2.add_run(); r2.text = body
        r2.font.size = Pt(12.5); r2.font.color.rgb = CARD_TEXT
        r2.font.name = "Noto Sans TC"

    bottom_y = row2_y + ch + Inches(0.18)
    txtbox(slide, ML, bottom_y, CW, Inches(0.35),
           "幾個沒有標準答案的問題，帶著走——",
           size=Pt(13), color=MUTED, align=PP_ALIGN.CENTER)
    txtbox(slide, ML, bottom_y + Inches(0.36), CW, Inches(0.46),
           "你願意為自己設定一個「不用 AI」的時段嗎？  ·  你的邊界，是工具決定的，還是你自己決定的？",
           size=Pt(14.5), color=WHITE, align=PP_ALIGN.CENTER)


# ── Main ─────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()
    s01_cover(prs)
    s02_problem(prs)
    s03_scenario(prs)
    s04_commands(prs)
    s05_case1(prs)
    s06_case2(prs)
    s07_case3(prs)
    s08_cleanlog(prs)
    s08b_playwright(prs)
    s09_tests(prs)
    s10_perf(prs)
    s11_correctness(prs)
    s12_comparison(prs)
    s13_why(prs)
    s14_quickstart(prs)
    s14b_mcp(prs)
    s14c_webui(prs)
    s15_bigquestion(prs)

    out = "/home/null/kw-filter/kw-filter.pptx"
    prs.save(out)
    print(f"Saved → {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
